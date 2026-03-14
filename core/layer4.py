"""Layer 4 — Query-driven evidence retrieval and answer synthesis.

Given a user question, Layer 4:
1. Uses LLM semantic judgement to binary-screen candidate files
2. Selects the most relevant files for deeper reading
3. Re-reads the most relevant files from disk
4. Produces an answer with supporting files and evidence
"""

from __future__ import annotations

import json
import math
import re
from pathlib import Path
from typing import Any
from collections import Counter

import pandas as pd
from openai import OpenAI

from .layer1 import layer1_technical

_STOPWORDS = {
    "the", "and", "for", "with", "that", "this", "from", "into", "about", "where",
    "what", "which", "when", "does", "have", "there", "their", "your", "our", "should",
    "would", "could", "using", "used", "use", "find", "look", "need", "question",
    "project", "file", "files", "document", "documents", "data", "layer", "query",
}


SEARCH_PLAN_PROMPT = """\
You help search an urban-design project archive.
Given a user question, produce a compact search plan for retrieving relevant files.

Return ONLY JSON with this shape:
{
  "intent": "<short sentence>",
  "search_terms": ["<keywords and phrases, bilingual if useful>"],
  "domain_hints": ["<likely domain names>"],
  "asset_hints": ["<likely asset types>"],
  "lifecycle_hints": ["<likely lifecycle stages>"],
  "answer_type": "fact|location|comparison|recommendation|list|unknown"
}

Rules:
- Keep search_terms concise and retrieval-oriented.
- CRITICAL: The archive files are in Danish and English. If the question is in
  Chinese or any other language, you MUST translate every key concept into English
  (and include relevant Danish equivalents if you know them) and add those
  translations as additional entries in search_terms. Never leave search_terms
  in only the question's language if that language is not Danish or English.
  Example: Chinese question about timber → add "timber", "wood", "tr\xe6", "construction".
- Domain hints must be from likely archive taxonomy such as Landscape & Public Realm,
  Urban Planning & Massing, Architecture & Buildings, Environment & Climate,
  Mobility & Transport, Administrative & Legal, Project Management, Reference & Research.
- Asset hints should be values like Data, Document, Drawing, Media, Archive.
- lifecycle_hints can be empty.
- If unsure, return broad but useful hints instead of inventing specifics.

Question: {question}
"""


SELECT_CANDIDATES_PROMPT = """\
You help shortlist archive files for deep reading.
You are given a user question and a list of candidate files with metadata and short summaries.

Return ONLY JSON with this shape:
{
    "selected_file_paths": ["<full path>", "<full path>"],
    "selection_reason": "<one sentence>",
    "notes": "<optional note about uncertainty>"
}

Rules:
- Select 2 to 5 file paths that are most worth reading in full.
- Prefer files whose summary, filename, domain, lifecycle, and path suggest direct evidence.
- If the question is about location, prioritize plans, reports, and summaries likely to mention places.
- If the question is about suitability or recommendation, prioritize specs, reports, data, and studies.
- Do not invent file paths.

Question:
{question}

Candidate previews:
{candidate_preview}
"""


BINARY_RELEVANCE_PROMPT = """\
You are screening archive files for a user question.
Decide relevance using semantic understanding, not keyword overlap.

Return ONLY JSON with this shape:
{
    "decisions": [
        {
            "file_path": "<full path exactly as provided>",
            "relevant": true,
            "confidence": 0.0,
            "reason": "<short reason>"
        }
    ]
}

Rules:
- Evaluate each file independently by meaning and likely evidence value.
- Mark relevant=true only if the file is likely useful for answering the question.
- confidence is 0.0 to 1.0.
- Keep reason concise (max 15 words).
- Do not invent file paths.

Question:
{question}

Files:
{batch_preview}
"""


ANSWER_PROMPT = """\
You answer questions about a design-project archive using ONLY the provided candidate files.
If the evidence is insufficient, say so clearly.

Return ONLY JSON with this shape:
{
  "answer": "<clear answer in the user's language if possible>",
  "confidence": "High|Medium|Low",
  "gaps": "<what is still uncertain or missing>",
  "relevant_files": [
    {
      "filename": "<file name>",
      "file_path": "<full path>",
      "why_it_matters": "<one sentence>",
      "evidence": "<short quote or extracted snippet>"
    }
  ]
}

Rules:
- Base the answer ONLY on the candidate files below.
- Prefer direct evidence over speculation.
- If the question asks "where", extract place/location clues.
- If the question asks for suitability/recommendation, explain what the documents support, and where evidence is weak.
- Keep the answer concise but useful.
- Include 2-6 relevant files when available.

User question:
{question}

Candidate files:
{candidate_context}
"""


def _safe_json_loads(raw: str, fallback: dict) -> dict:
    try:
        cleaned = re.sub(r"^```(?:json)?\s*", "", raw.strip(), flags=re.IGNORECASE)
        cleaned = re.sub(r"\s*```$", "", cleaned, flags=re.IGNORECASE)
        return json.loads(cleaned)
    except Exception:
        return fallback


def _tokenize(text: str) -> list[str]:
    text = (text or "").lower()
    latin_tokens = re.findall(r"[a-z0-9][a-z0-9\-_]{1,}", text)
    phrase_tokens = re.findall(r"[\u4e00-\u9fff]{2,}", text)
    tokens = []
    for token in latin_tokens + phrase_tokens:
        if token not in _STOPWORDS and token not in tokens:
            tokens.append(token)
    return tokens


def build_search_plan(question: str, client: OpenAI, model: str) -> dict:
    """Use the LLM to derive retrieval hints from the user question."""
    fallback = {
        "intent": question.strip(),
        "search_terms": _tokenize(question)[:8],
        "domain_hints": [],
        "asset_hints": [],
        "lifecycle_hints": [],
        "answer_type": "unknown",
    }

    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Respond only with valid JSON."},
                {"role": "user", "content": SEARCH_PLAN_PROMPT.format(question=question.strip())},
            ],
            temperature=0,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content or ""
        plan = _safe_json_loads(raw, fallback)
    except Exception:
        plan = fallback

    plan.setdefault("intent", question.strip())
    plan.setdefault("search_terms", fallback["search_terms"])
    plan.setdefault("domain_hints", [])
    plan.setdefault("asset_hints", [])
    plan.setdefault("lifecycle_hints", [])
    plan.setdefault("answer_type", "unknown")
    return plan


def _score_text_match(text: str, search_terms: list[str], weight: float) -> float:
    # kept for any external callers; not used in ranking anymore
    hay = (text or "").lower()
    return sum(weight for t in search_terms if str(t or "").strip().lower() in hay)


# ── BM25 retrieval ─────────────────────────────────────────────────────────
_BM25_FIELDS = [
    "filename", "keywords", "short_summary", "_reasoning",
    "review_reasons", "information_type", "domain",
    "lifecycle", "asset_type", "file_path",
]


def _build_bm25_corpus(df: pd.DataFrame) -> tuple[list[list[str]], float, Counter, int]:
    """Tokenise every row into a combined doc and compute BM25 corpus stats."""
    docs: list[list[str]] = []
    for _, row in df.iterrows():
        combined = " ".join(str(row.get(col, "")) for col in _BM25_FIELDS)
        docs.append(_tokenize(combined))
    N = len(docs)
    df_freq: Counter = Counter()
    for doc in docs:
        df_freq.update(set(doc))
    avg_dl = sum(len(d) for d in docs) / max(N, 1)
    return docs, avg_dl, df_freq, N


def _bm25_score(
    query_tokens: list[str],
    doc_tokens: list[str],
    avg_dl: float,
    df_freq: Counter,
    N: int,
    k1: float = 1.5,
    b: float = 0.75,
) -> float:
    """BM25 score for one document against query tokens."""
    if not query_tokens or not doc_tokens:
        return 0.0
    dl = len(doc_tokens)
    tf_counts = Counter(doc_tokens)
    score = 0.0
    for term in query_tokens:
        tf = tf_counts.get(term, 0)
        if tf == 0:
            continue
        n_t = df_freq.get(term, 0)
        idf = math.log((N - n_t + 0.5) / (n_t + 0.5) + 1.0)
        tf_norm = (tf * (k1 + 1)) / (tf + k1 * (1 - b + b * dl / max(avg_dl, 1)))
        score += idf * tf_norm
    return score


def _build_binary_batch_preview(batch: pd.DataFrame) -> str:
    parts: list[str] = []
    for idx, (_, row) in enumerate(batch.iterrows(), start=1):
        parts.append(
            f"[File {idx}]\n"
            f"Filename: {row.get('filename', '')}\n"
            f"Path: {row.get('file_path', '')}\n"
            f"Domain: {row.get('domain', '')}\n"
            f"Asset Type: {row.get('asset_type', '')}\n"
            f"Lifecycle: {row.get('lifecycle', '')}\n"
            f"Summary: {str(row.get('short_summary', ''))[:220]}\n"
            f"Reasoning: {str(row.get('_reasoning', ''))[:220]}\n"
        )
    return "\n\n".join(parts)


def rank_candidate_rows(
    df: pd.DataFrame,
    question: str,
    client: OpenAI,
    model: str,
    top_k: int = 12,
    batch_size: int = 20,
) -> pd.DataFrame:
    """LLM-first binary relevance filter over processed rows (semantic, not keyword-based)."""
    if df.empty:
        return df.copy()

    ranked = df.copy().reset_index(drop=True)
    relevance_score = pd.Series([0.0] * len(ranked), index=ranked.index, dtype="float")
    is_relevant = pd.Series([False] * len(ranked), index=ranked.index, dtype="bool")

    for start in range(0, len(ranked), max(1, batch_size)):
        batch = ranked.iloc[start:start + max(1, batch_size)]
        try:
            resp = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "Respond only with valid JSON."},
                    {
                        "role": "user",
                        "content": BINARY_RELEVANCE_PROMPT.format(
                            question=question.strip(),
                            batch_preview=_build_binary_batch_preview(batch),
                        ),
                    },
                ],
                temperature=0,
                response_format={"type": "json_object"},
            )
            raw = resp.choices[0].message.content or ""
            parsed = _safe_json_loads(raw, {})
            decisions = parsed.get("decisions", [])

            by_path = {
                str(item.get("file_path", "")).strip(): item
                for item in decisions
                if str(item.get("file_path", "")).strip()
            }
            for idx, row in batch.iterrows():
                path_key = str(row.get("file_path", "")).strip()
                dec = by_path.get(path_key)
                if not dec:
                    continue
                rel = bool(dec.get("relevant", False))
                conf_raw = dec.get("confidence", 0.5)
                try:
                    conf = max(0.0, min(1.0, float(conf_raw)))
                except Exception:
                    conf = 0.5
                is_relevant.loc[idx] = rel
                relevance_score.loc[idx] = conf if rel else 0.0
        except Exception:
            continue

    ranked["_query_score"] = relevance_score
    relevant_rows = ranked[is_relevant].copy()
    if not relevant_rows.empty:
        relevant_rows = relevant_rows.sort_values(
            ["_query_score", "confidence", "year"],
            ascending=[False, False, False],
            na_position="last",
        )
        return relevant_rows.head(top_k)

    # Fallback: no clear relevant rows returned; keep best-confidence rows for next LLM stage
    ranked = ranked.sort_values(["confidence", "year"], ascending=[False, False], na_position="last")
    ranked["_query_score"] = 0.0
    return ranked.head(min(top_k, len(ranked)))


def _read_full_file_content(file_path: Path, max_chars: int = 24000) -> tuple[str, str]:
    """Read a deeper/full content slice for query answering, beyond layer1 short samples."""
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            parts: list[str] = []
            for page_idx, page in enumerate(reader.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    parts.append(f"[p{page_idx + 1}] {text}")
                if sum(len(p) for p in parts) >= max_chars:
                    break
            raw = "\n".join(parts)
            return raw[:max_chars] or "[PDF has little or no extractable text]", f"fuller PDF pass across {min(len(reader.pages), page_idx + 1 if 'page_idx' in locals() else 0)} page(s)"

        if ext in {".docx", ".doc"}:
            import docx
            doc = docx.Document(str(file_path))
            parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
                    if sum(len(p) for p in parts) >= max_chars:
                        break
            raw = "\n".join(parts)
            return raw[:max_chars] or "[Document has little extractable text]", "fuller DOCX pass"

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides):
                texts = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame and shape.text.strip()]
                if texts:
                    parts.append(f"[Slide {i+1}] " + " | ".join(texts))
                if sum(len(p) for p in parts) >= max_chars:
                    break
            raw = "\n".join(parts)
            return raw[:max_chars] or "[Presentation has little extractable text]", "fuller PPTX pass"

        if ext in {".txt", ".json"}:
            raw = file_path.read_text(encoding="utf-8", errors="replace")
            return raw[:max_chars], "full text read"

        if ext == ".csv":
            df = pd.read_csv(file_path, encoding="utf-8", errors="replace")
            raw = f"Columns: {list(df.columns)}\n{df.head(200).to_string()}"
            return raw[:max_chars], "CSV deeper tabular read"

        if ext in {".xlsx", ".xls"}:
            xl = pd.ExcelFile(str(file_path))
            blocks: list[str] = [f"Sheets: {xl.sheet_names}"]
            for sheet in xl.sheet_names[:5]:
                df = pd.read_excel(file_path, sheet_name=sheet, nrows=80)
                blocks.append(f"[Sheet: {sheet}]\nColumns: {list(df.columns)}\n{df.to_string()}")
                if sum(len(b) for b in blocks) >= max_chars:
                    break
            raw = "\n\n".join(blocks)
            return raw[:max_chars], "Excel deeper read"

        if ext == ".eml":
            import email
            from email import policy as email_policy
            with open(file_path, encoding="utf-8", errors="replace") as f:
                msg = email.message_from_file(f, policy=email_policy.default)
            body_parts = []
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body_parts.append(part.get_content() or "")
            else:
                body_parts.append(msg.get_content() or "")
            raw = f"From: {msg.get('From', '')}\nSubject: {msg.get('Subject', '')}\nDate: {msg.get('Date', '')}\n\n" + "\n".join(body_parts)
            return raw[:max_chars], "fuller email read"
    except Exception:
        pass

    fallback = layer1_technical(file_path)
    return str(fallback.get("content_sample", ""))[:max_chars], str(fallback.get("extraction_coverage", "layer1 fallback"))


def _build_candidate_preview(candidates: pd.DataFrame) -> str:
    parts: list[str] = []
    for idx, (_, row) in enumerate(candidates.iterrows(), start=1):
        parts.append(
            f"[Preview {idx}]\n"
            f"Filename: {row.get('filename', '')}\n"
            f"Path: {row.get('file_path', '')}\n"
            f"Domain: {row.get('domain', '')}\n"
            f"Asset Type: {row.get('asset_type', '')}\n"
            f"Lifecycle: {row.get('lifecycle', '')}\n"
            f"Score: {float(row.get('_query_score', 0)):.1f}\n"
            f"Summary: {row.get('short_summary', '')}\n"
            f"Reasons: {row.get('review_reasons', '')}\n"
        )
    return "\n\n".join(parts)


def select_files_for_deep_read(question: str, candidates: pd.DataFrame, client: OpenAI, model: str, max_select: int = 4) -> pd.DataFrame:
    """Use the LLM to choose which candidate files deserve deeper reading."""
    if candidates.empty:
        return candidates

    fallback = candidates.head(max_select)
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Respond only with valid JSON."},
                {
                    "role": "user",
                    "content": SELECT_CANDIDATES_PROMPT.format(
                        question=question.strip(),
                        candidate_preview=_build_candidate_preview(candidates),
                    ),
                },
            ],
            temperature=0,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content or ""
        parsed = _safe_json_loads(raw, {})
        selected_paths = [str(x) for x in parsed.get("selected_file_paths", []) if str(x).strip()]
        if selected_paths:
            selected = candidates[candidates["file_path"].astype(str).isin(selected_paths)]
            if not selected.empty:
                return selected.head(max_select)
    except Exception:
        pass

    return fallback


def reread_candidate_files(candidates: pd.DataFrame, max_files: int = 6) -> list[dict[str, Any]]:
    """Re-read selected files from disk using deeper/full-content reads when possible."""
    reread_docs: list[dict[str, Any]] = []

    for _, row in candidates.head(max_files).iterrows():
        file_path = Path(str(row.get("file_path", ""))).expanduser()
        if not file_path.exists() or not file_path.is_file():
            reread_docs.append({
                "filename": str(row.get("filename", file_path.name or "Unknown")),
                "file_path": str(file_path),
                "domain": row.get("domain", "Unknown"),
                "asset_type": row.get("asset_type", "Unknown"),
                "lifecycle": row.get("lifecycle", "Unknown"),
                "summary": row.get("short_summary", ""),
                "coverage": "file missing at query time",
                "content_sample": "[File not found on disk during query reread]",
                "query_score": float(row.get("_query_score", 0)),
            })
            continue

        try:
            content_sample, coverage = _read_full_file_content(file_path)
            reread = layer1_technical(file_path)
            reread_docs.append({
                "filename": reread.get("filename", file_path.name),
                "file_path": reread.get("file_path", str(file_path)),
                "domain": row.get("domain", "Unknown"),
                "asset_type": row.get("asset_type", "Unknown"),
                "lifecycle": row.get("lifecycle", "Unknown"),
                "summary": row.get("short_summary", ""),
                "coverage": coverage or reread.get("extraction_coverage", ""),
                "content_sample": content_sample[:24000],
                "query_score": float(row.get("_query_score", 0)),
            })
        except Exception as exc:
            reread_docs.append({
                "filename": str(row.get("filename", file_path.name)),
                "file_path": str(file_path),
                "domain": row.get("domain", "Unknown"),
                "asset_type": row.get("asset_type", "Unknown"),
                "lifecycle": row.get("lifecycle", "Unknown"),
                "summary": row.get("short_summary", ""),
                "coverage": f"reread failed: {exc}",
                "content_sample": "[Could not re-read file content]",
                "query_score": float(row.get("_query_score", 0)),
            })

    return reread_docs


def _build_candidate_context(docs: list[dict[str, Any]]) -> str:
    parts: list[str] = []
    for idx, doc in enumerate(docs, start=1):
        parts.append(
            f"[Candidate {idx}]\n"
            f"Filename: {doc.get('filename', '')}\n"
            f"Path: {doc.get('file_path', '')}\n"
            f"Domain: {doc.get('domain', '')}\n"
            f"Asset Type: {doc.get('asset_type', '')}\n"
            f"Lifecycle: {doc.get('lifecycle', '')}\n"
            f"Query Score: {doc.get('query_score', 0):.1f}\n"
            f"Pipeline Summary: {doc.get('summary', '')}\n"
            f"Extraction Coverage: {doc.get('coverage', '')}\n"
            f"Content Sample:\n{doc.get('content_sample', '')}\n"
        )
    return "\n\n".join(parts)


def synthesize_query_answer(question: str, reread_docs: list[dict[str, Any]], client: OpenAI, model: str) -> dict:
    """Ask the LLM to answer using only the re-read candidate documents."""
    fallback_files = [
        {
            "filename": doc.get("filename", ""),
            "file_path": doc.get("file_path", ""),
            "why_it_matters": doc.get("summary", "Potentially relevant file."),
            "evidence": str(doc.get("content_sample", ""))[:200],
        }
        for doc in reread_docs[:3]
    ]
    fallback = {
        "answer": "I could not generate a reliable answer from the available files.",
        "confidence": "Low",
        "gaps": "Layer 4 synthesis failed or the evidence was too weak.",
        "relevant_files": fallback_files,
    }

    if not reread_docs:
        return {
            "answer": "No candidate files were available for this question.",
            "confidence": "Low",
            "gaps": "No files matched the query strongly enough to analyze.",
            "relevant_files": [],
        }

    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Respond only with valid JSON."},
                {
                    "role": "user",
                    "content": ANSWER_PROMPT.format(
                        question=question.strip(),
                        candidate_context=_build_candidate_context(reread_docs),
                    ),
                },
            ],
            temperature=0,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content or ""
        answer = _safe_json_loads(raw, fallback)
    except Exception:
        answer = fallback

    answer.setdefault("answer", fallback["answer"])
    answer.setdefault("confidence", "Low")
    answer.setdefault("gaps", "")
    answer.setdefault("relevant_files", fallback_files)
    return answer


def layer4_query(
    question: str,
    processed_df: pd.DataFrame,
    client: OpenAI,
    model: str,
    top_k: int = 8,
    reread_k: int = 6,
) -> dict:
    """End-to-end Layer 4 question answering over processed archive results."""
    if processed_df is None or processed_df.empty:
        return {
            "question": question,
            "search_plan": {},
            "answer": "No processed files are available for querying.",
            "confidence": "Low",
            "gaps": "Run the pipeline first so Layer 4 has files to search.",
            "relevant_files": [],
            "candidate_count": 0,
        }

    plan = {
        "mode": "llm_binary_relevance",
        "intent": question.strip(),
        "search_terms": [],
        "domain_hints": [],
        "asset_hints": [],
        "lifecycle_hints": [],
        "answer_type": "unknown",
    }
    candidates = rank_candidate_rows(
        processed_df,
        question=question,
        client=client,
        model=model,
        top_k=max(top_k, 10),
    )
    deep_read_candidates = select_files_for_deep_read(question, candidates, client=client, model=model, max_select=reread_k)
    reread_docs = reread_candidate_files(deep_read_candidates, max_files=reread_k)
    answer = synthesize_query_answer(question, reread_docs, client=client, model=model)

    return {
        "question": question,
        "search_plan": plan,
        "answer": answer.get("answer", ""),
        "confidence": answer.get("confidence", "Low"),
        "gaps": answer.get("gaps", ""),
        "relevant_files": answer.get("relevant_files", []),
        "candidate_count": int(len(candidates)),
        "deep_read_count": int(len(deep_read_candidates)),
        "candidates": candidates[[c for c in [
            "filename", "file_path", "domain", "asset_type", "lifecycle", "short_summary", "_query_score"
        ] if c in candidates.columns]].to_dict("records"),
    }
