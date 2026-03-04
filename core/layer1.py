import re
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Optional

# ── Info type map ─────────────────────────────────────────────────────────
INFO_TYPE_MAP = {
    ".pdf":     "Document — see content_sample",
    ".docx":    "Narrative / Textual",
    ".doc":     "Narrative / Textual",
    ".txt":     "Narrative / Textual",
    ".pptx":    "Visual / Presentation",
    ".xlsx":    "Quantitative / Tabular",
    ".xls":     "Quantitative / Tabular",
    ".csv":     "Quantitative / Tabular",
    ".json":    "Quantitative / Tabular",
    ".dwg":     "Schematic / Technical",
    ".dxf":     "Schematic / Technical",
    ".dwf":     "Schematic / Technical",          # ← NEW: Design Web Format
    ".ifc":     "Schematic / BIM",
    ".rvt":     "Schematic / BIM",
    ".nwd":     "Schematic / BIM",                # ← NEW: Navisworks
    ".jpg":     "Visual / Media",
    ".jpeg":    "Visual / Media",
    ".png":     "Visual / Media",
    ".tiff":    "Visual / Media",
    ".tif":     "Visual / Media",
    ".mp4":     "Visual / Media",                 # ← NEW: video
    ".mov":     "Visual / Media",                 # ← NEW: video
    ".shp":     "Spatial / Cartographic",
    ".geojson": "Spatial / Cartographic",
    ".kml":     "Spatial / Cartographic",
    ".gpkg":    "Spatial / Cartographic",         # ← NEW: GeoPackage
    ".eml":     "Narrative / Email",
    ".msg":     "Narrative / Email",
    ".zip":     "Archive",
    ".7z":      "Archive",                        # ← NEW
    ".rar":     "Archive",                        # ← NEW
}

# ── PDF extraction config ─────────────────────────────────────────────────
PDF_HEAD_PAGES = 3
PDF_TAIL_PAGES = 2
PDF_MID_SAMPLE = 2
PDF_MAX_CHARS  = 2400
DEFAULT_MAX_CHARS = 800

# Data-detection moved to Layer 2 to centralise classification logic.


# Size thresholds (KB)
_SIZE_CATEGORIES = [
    (0,     50,    "tiny"),
    (50,    500,   "small"),
    (500,   5_000, "medium"),
    (5_000, 50_000, "large"),
    (50_000, None,  "xlarge"),
]


# ── Year extraction ───────────────────────────────────────────────────────
def _find_years_in_text(text: str) -> list[str]:
    current = datetime.now().year
    hits    = re.findall(r"\b(20\d{2}|19\d{2})\b", text)
    return [h for h in hits if int(h) <= current]


def _extract_year(filename: str, content: str = "", file_path: Optional[Path] = None) -> Optional[str]:
    """
    Priority order:
      1. Year in filename (most reliable)
      2. Most-frequent year in content (noise-filtered)
      3. File OS modification date (fallback — better than nothing)
    """
    filename_years = _find_years_in_text(filename)
    if filename_years:
        return max(filename_years, key=int)

    if content:
        cleaned = re.sub(r"\b[A-Z]{1,5}[\s\-]\d{3,6}[-:]\d{2,4}\b", "", content)
        cleaned = re.sub(r"©\s*\d{4}", "", cleaned)
        content_years = _find_years_in_text(cleaned)
        if content_years:
            freq       = Counter(content_years)
            max_freq   = max(freq.values())
            candidates = [y for y, c in freq.items() if c == max_freq]
            return max(candidates, key=int)

    # NEW ── fallback: OS modification date
    if file_path:
        try:
            mtime = file_path.stat().st_mtime
            mtime_year = str(datetime.fromtimestamp(mtime).year)
            current = datetime.now().year
            if 2000 <= int(mtime_year) <= current:
                return mtime_year + " (mtime)"
        except Exception:
            pass

    return None


# ── Coverage helpers ──────────────────────────────────────────────────────
def _cov_pages(read: list[int], total: int) -> str:
    return f"{len(read)}/{total} pages sampled"

def _cov_chars(content: str, cap: int, label: str = "") -> str:
    truncated = len(content) >= cap
    base = f"truncated at {cap} chars" if truncated else f"full content ({len(content)} chars)"
    return f"{label} — {base}" if label else base


# ── Data hint — content structure analysis ────────────────────────────────
def _score_content_structure(content: str) -> int:
    """
    Analyse the structural pattern of extracted text to decide if it looks
    like tabular/structured data.  Language-agnostic: works in Arabic, French,
    Chinese, or any naming convention because it measures *shape*, not words.

    Signals measured
    ────────────────
    numeric_density   — what fraction of tokens are numbers?
                        Data files are full of numbers; narrative docs are not.
    delimiter_density — frequency of field separators (, | \t ;) per line.
                        High density → columnar / tabular structure.
    row_regularity    — do lines have a similar token count AND high delimiter
                        density?  Checked together because word-wrapped prose
                        also has regular line lengths (false positive if checked
                        alone).
    header_pattern    — does the first non-empty line look like a header row
                        (3+ delimiter-separated short tokens, matching count
                        in the second line)?

    Each signal contributes ±points independently so they combine gracefully.
    No language keywords are used.
    """
    if not content or len(content.strip()) < 40:
        return 0

    lines = [l for l in content.splitlines() if l.strip()]
    if len(lines) < 2:
        return 0

    score = 0

    # ── Numeric density ───────────────────────────────────────────────────
    tokens = re.split(r"[\s,|\t;/]+", content)
    tokens = [t for t in tokens if t]
    if tokens:
        numeric = sum(1 for t in tokens if re.fullmatch(r"-?\d+([.,]\d+)?", t))
        num_ratio = numeric / len(tokens)
        if   num_ratio > 0.35:  score += 3
        elif num_ratio > 0.15:  score += 1
        elif num_ratio < 0.03:  score -= 1

    # ── Delimiter density ─────────────────────────────────────────────────
    delimiters = re.findall(r"[,|\t;]", content)
    delim_per_line = len(delimiters) / len(lines)
    if   delim_per_line > 4:   score += 3
    elif delim_per_line > 1.5: score += 1
    elif delim_per_line < 0.3: score -= 1

    # ── Row regularity (ONLY when delimiters also present) ────────────────
    # Word-wrapped prose also has regular line lengths → checking regularity
    # alone causes false positives.  We require delim_per_line > 1 as a gate
    # so this bonus only fires for genuinely columnar content.
    if delim_per_line > 1.0 and len(lines) >= 3:
        lengths  = [len(l.split()) for l in lines]
        mean_len = sum(lengths) / len(lengths)
        if mean_len > 0:
            variance = sum((x - mean_len) ** 2 for x in lengths) / len(lengths)
            cv       = (variance ** 0.5) / mean_len
            if   cv < 0.3:  score += 2
            elif cv < 0.6:  score += 1

    # ── Header pattern ────────────────────────────────────────────────────
    first     = lines[0]
    cols_first = re.split(r"[,|\t;]", first)
    if len(cols_first) >= 3:
        avg_col_len = sum(len(c.strip()) for c in cols_first) / len(cols_first)
        if avg_col_len < 30:
            score += 2
            if len(lines) > 1:
                cols_second = re.split(r"[,|\t;]", lines[1])
                if len(cols_second) == len(cols_first):
                    score += 1

    return score


def _is_data_hint(file_path: Path, content: str) -> str:
    """Placeholder kept for backward compatibility.

    The real data-detection logic has been moved to Layer 2. Layer 1 will
    no longer assert data-vs-document; Layer 2 computes `is_data_hint`
    from `file_path` and `content_sample` when classifying.
    """
    return "Unlikely"


# ── NEW: Path segments ────────────────────────────────────────────────────
def _extract_path_segments(file_path: Path) -> list[str]:
    """
    Returns all meaningful folder names in the path (from root down to parent),
    stripping drive letters and common no-signal roots like 'Users', 'home', etc.

    Example:
      /Users/luca/Projects/DataTaxonomy/2022Masterplan/03_Design/ARCH/drawings/file.dwg
      → ['DataTaxonomy', '2022Masterplan', '03_Design', 'ARCH', 'drawings']

    This is far more informative than just `folder = file_path.parent.name`
    and gives Layer 2 the full context chain.
    """
    _SKIP = {"users", "user", "home", "documents", "downloads", "desktop",
             "onedrive", "sharepoint", "sites", "shared documents",
             "c:", "d:", "volumes", "mnt", "/"}
    parts = []
    for part in file_path.parts[:-1]:  # exclude filename
        clean = part.strip("/\\")
        if clean.lower() not in _SKIP and clean not in ("", "."):
            parts.append(clean)
    return parts[-8:]  # cap at 8 levels — enough context, not noise


# ── Filename signals — structural token extraction ────────────────────────
def _extract_filename_signals(filename: str) -> dict:
    """
    Structural parsing of filename segments.  No enumeration of discipline
    codes or status tags — those lists are English-only and naming-convention-
    specific.  Instead we extract token shapes and let Layer 2 (LLM) interpret
    the meaning.

    What we extract
    ───────────────
    code_tokens      — all-caps tokens, 2–6 chars, no digits.
                       These are positional abbreviations (discipline, status,
                       zone, phase codes).  Passed raw to LLM; it understands
                       ARCH, MEP, IFC, WIP, GIS regardless of project language.

    version          — structural version/revision pattern: V/R/P + number,
                       or bare two-digit sequence between separators.
                       Language-agnostic (same regex works in any project).

    drawing_number   — letter-code + separator + 3–5 digit run (A-001, SK-024).
                       Strong structural signal that the file is a drawing.

    has_date_in_name — True if YYYYMMDD or YYYY-MM-DD appears in stem.

    numeric_tokens   — standalone numbers between separators (zone IDs, rev
                       numbers, phase indices like 03, 204, 1001).

    token_count      — total segment count after splitting on separators.
                       Very high count (>6) → likely structured naming convention.
                       Very low count (1–2) → freeform or legacy name.

    raw_stem         — full stem passed to Layer 2 so the LLM can read it
                       directly without lossy pre-parsing.
    """
    stem       = Path(filename).stem
    stem_upper = stem.upper()

    # Split on all common separators
    tokens = re.split(r"[_\-\.\s]+", stem_upper)
    tokens = [t for t in tokens if t]

    # All-caps code tokens: 2–6 letters, no digits — discipline / status / zone codes
    code_tokens = [t for t in tokens if re.fullmatch(r"[A-Z]{2,6}", t)]

    # Version/revision — language-agnostic structural pattern
    version_match = re.search(
        r"[\-_](?:V|REV|R|P)[\-_]?(\d{1,3}|[A-F])(?=[\-_.]|$)"
        r"|[\-_](\d{2})(?=[\-_.]|$)",
        stem_upper,
    )
    version = version_match.group(0).strip("-_") if version_match else None

    # Drawing number: letter prefix + separator + 3–5 digit run
    drawing_match  = re.search(r"\b([A-Z]{1,4}[\-\.]\d{3,5})\b", stem_upper)
    drawing_number = drawing_match.group(1) if drawing_match else None

    # Date in filename (YYYYMMDD compact or YYYY-MM-DD / YYYY_MM_DD)
    has_date_in_name = bool(
        re.search(r"20\d{2}[0-1]\d[0-3]\d", filename) or
        re.search(r"20\d{2}[-_][0-1]\d[-_][0-3]\d", filename)
    )

    # Standalone numeric segments (zone, phase, sequence IDs)
    numeric_tokens = [t for t in tokens if re.fullmatch(r"\d{1,5}", t)]

    return {
        "code_tokens":       code_tokens,        # raw; LLM interprets meaning
        "version":           version,
        "drawing_number":    drawing_number,
        "has_date_in_name":  has_date_in_name,
        "numeric_tokens":    numeric_tokens,
        "token_count":       len(tokens),
        "raw_stem":          stem,               # full stem for LLM
    }


# ── NEW: Size category ────────────────────────────────────────────────────
def _size_category(size_kb: int) -> str:
    for lo, hi, label in _SIZE_CATEGORIES:
        if hi is None or size_kb < hi:
            if size_kb >= lo:
                return label
    return "xlarge"


# ── Content extraction ────────────────────────────────────────────────────
def _extract_content(file_path: Path) -> tuple[str, Optional[int], str]:
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            r     = PdfReader(str(file_path))
            pages = len(r.pages)
            head  = list(range(min(PDF_HEAD_PAGES, pages)))
            tail  = list(range(max(0, pages - PDF_TAIL_PAGES), pages))
            mid   = []
            if pages > PDF_HEAD_PAGES + PDF_TAIL_PAGES + 2 and PDF_MID_SAMPLE > 0:
                mid_start = PDF_HEAD_PAGES
                mid_end   = pages - PDF_TAIL_PAGES
                step      = max(1, (mid_end - mid_start) // (PDF_MID_SAMPLE + 1))
                mid = [mid_start + step * (i + 1) for i in range(PDF_MID_SAMPLE)
                       if mid_start + step * (i + 1) < mid_end]
            page_indices = list(dict.fromkeys(head + mid + tail))
            parts = []
            for i in page_indices:
                text = (r.pages[i].extract_text() or "").strip()
                if text:
                    parts.append(f"[p{i + 1}] {text}")
            if parts:
                raw      = "\n".join(parts)
                content  = raw[:PDF_MAX_CHARS]
                coverage = _cov_pages(page_indices, pages)
                if len(raw) > PDF_MAX_CHARS:
                    coverage += f" — truncated at {PDF_MAX_CHARS} chars"
            else:
                size_kb  = file_path.stat().st_size // 1024
                content  = (f"[PDF: no extractable text — likely scanned drawing or print export]\n"
                            f"Pages: {pages} | Size: {size_kb} KB\n"
                            f"Infer content from filename and folder path.")
                coverage = f"{pages} pages — no text layer detected"
        except Exception as e:
            content, pages, coverage = f"[PDF error: {e}]", None, "extraction failed"
        return content, pages, coverage

    elif ext in (".docx", ".doc"):
        try:
            import docx
            doc        = docx.Document(str(file_path))
            para_texts = [p.text for p in doc.paragraphs[:20] if p.text.strip()]
            table_texts = []
            for tbl_idx, table in enumerate(doc.tables[:3]):
                table_texts.append(f"[Table {tbl_idx + 1}]")
                for row in table.rows[:3]:
                    row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if row_text:
                        table_texts.append(row_text)
            parts    = para_texts + (["[Tables:]"] + table_texts if table_texts else [])
            content  = "\n".join(parts)[:DEFAULT_MAX_CHARS]
            coverage = (f"{len(para_texts)} paragraphs"
                        + (f" + {len(doc.tables)} table(s) sampled" if doc.tables else ""))
        except Exception as e:
            content  = f"[DOCX error — likely legacy .doc binary: {e}]"
            coverage = "extraction failed"
        return content, None, coverage

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs        = Presentation(str(file_path))
            all_slides = list(prs.slides)
            sample     = all_slides[:5]
            parts      = []
            for i, slide in enumerate(sample):
                texts = [shape.text.strip() for shape in slide.shapes
                         if shape.has_text_frame and shape.text.strip()]
                if texts:
                    parts.append(f"[Slide {i+1}] " + " | ".join(texts[:3]))
            content  = "\n".join(parts)[:DEFAULT_MAX_CHARS]
            coverage = f"{len(sample)}/{len(all_slides)} slides sampled"
        except Exception as e:
            content, coverage = f"[PPTX error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext in (".xlsx", ".xls"):
        try:
            import pandas as pd
            xl      = pd.ExcelFile(str(file_path))
            df      = pd.read_excel(file_path, sheet_name=xl.sheet_names[0], nrows=5)
            headers = list(df.columns)
            content = (f"Sheets: {xl.sheet_names}\nColumns: {headers}\n"
                       f"{df.head(5).to_string()}")[:DEFAULT_MAX_CHARS]
            coverage = f"sheet 1 of {len(xl.sheet_names)}, first 5 rows + headers"
        except Exception as e:
            content, coverage = f"[Excel error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext == ".csv":
        try:
            import pandas as pd
            df      = pd.read_csv(file_path, nrows=5, encoding="utf-8", errors="replace")
            headers = list(df.columns)
            content = (f"Columns: {headers}\n{df.head(5).to_string()}")[:DEFAULT_MAX_CHARS]
            coverage = "first 5 rows + headers"
        except Exception as e:
            content, coverage = f"[CSV error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext in (".txt", ".json"):
        try:
            raw      = file_path.read_text(encoding="utf-8", errors="replace")
            content  = raw[:DEFAULT_MAX_CHARS]
            coverage = _cov_chars(raw, DEFAULT_MAX_CHARS)
        except Exception as e:
            content, coverage = f"[Text error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext == ".dwg":
        # NEW ── Try to read ASCII header bytes (DWG starts with "AC" version string)
        try:
            raw_bytes = file_path.read_bytes()[:200]
            header    = raw_bytes[:6].decode("ascii", errors="replace")
            # DWG version map
            _DWG_VER = {
                "AC1032": "AutoCAD 2018–2021", "AC1027": "AutoCAD 2013–2017",
                "AC1024": "AutoCAD 2010–2012", "AC1021": "AutoCAD 2007–2009",
                "AC1018": "AutoCAD 2004–2006", "AC1015": "AutoCAD 2000–2002",
            }
            ver_label = _DWG_VER.get(header.strip(), header.strip())
            content   = f"[DWG binary — version: {ver_label} — filename and folder path signals only]"
            coverage  = "binary header only"
        except Exception:
            content  = "[DWG binary — filename and folder path signals only]"
            coverage = "binary, no extraction"
        return content, None, coverage

    elif ext == ".dxf":
        try:
            raw       = file_path.read_text(encoding="utf-8", errors="replace")[:1200]
            layers    = re.findall(r"(?<=\n  2\n)[A-Z0-9_\-]+(?=\n)", raw)
            layer_str = ", ".join(dict.fromkeys(layers[:20]))
            content   = (f"[DXF layers: {layer_str}]\n{raw[:400]}" if layer_str else raw[:400])
            coverage  = f"{len(layers)} layer names extracted"
        except Exception as e:
            content, coverage = f"[DXF error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext == ".rvt":
        return "[Revit binary — filename and folder path signals only]", None, "binary, no extraction"

    elif ext == ".ifc":
        # NEW ── Structured IFC header + entity type frequency
        try:
            raw = file_path.read_text(encoding="utf-8", errors="replace")

            # Extract STEP header fields
            desc_match   = re.search(r"FILE_DESCRIPTION\s*\(([^;]+)\)", raw)
            schema_match = re.search(r"FILE_SCHEMA\s*\(\s*\('([^']+)'\)", raw)
            name_match   = re.search(r"FILE_NAME\s*\('([^']*)'", raw)

            description = desc_match.group(1)[:200].strip() if desc_match else ""
            schema      = schema_match.group(1) if schema_match else ""
            file_name   = name_match.group(1) if name_match else ""

            # Count IFC entity types (e.g. IFCWALL, IFCSLAB, IFCCOLUMN …)
            entity_hits = re.findall(r"^#\d+=\s*(IFC[A-Z0-9]+)\(", raw, re.MULTILINE)
            top_entities = ", ".join(
                f"{e}×{c}" for e, c in Counter(entity_hits).most_common(8)
            )

            header_block = "\n".join(filter(None, [
                f"[IFC schema: {schema}]" if schema else "",
                f"[Authored as: {file_name}]" if file_name else "",
                f"[Description: {description}]" if description else "",
                f"[Top entity types: {top_entities}]" if top_entities else "",
            ]))
            content  = (header_block + "\n" + raw[:800])[:DEFAULT_MAX_CHARS]
            coverage = _cov_chars(raw, DEFAULT_MAX_CHARS, "IFC header+entities")
        except Exception as e:
            content, coverage = f"[IFC error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".tif"):
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
            img           = Image.open(file_path)
            size_kb       = file_path.stat().st_size // 1024
            exif_year     = ""
            exif_date_raw = ""
            try:
                exif_data = img._getexif() or {}
                for tag_id, val in exif_data.items():
                    if TAGS.get(tag_id) == "DateTimeOriginal":
                        exif_date_raw = str(val)[:10]
                        exif_year     = exif_date_raw[:4]
                        break
            except Exception:
                pass
            exif_tag = f" exif_year={exif_year}" if exif_year else ""
            content  = f"[Image {img.width}x{img.height}px mode={img.mode} size={size_kb}KB{exif_tag}]"
            coverage = "EXIF metadata" + (f" — shoot date {exif_date_raw}" if exif_date_raw else " — no EXIF date")
        except Exception as e:
            content, coverage = f"[Image error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext == ".eml":
        try:
            import email
            from email import policy as email_policy
            with open(file_path, encoding="utf-8", errors="replace") as f:
                msg = email.message_from_file(f, policy=email_policy.default)
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = (part.get_content() or "")[:300]
                        break
            else:
                body = (msg.get_content() or "")[:300]
            content  = (f"From: {msg.get('From', '')}\nSubject: {msg.get('Subject', '')}\n"
                        f"Date: {msg.get('Date', '')}\nBody: {body.strip()}")
            coverage = "headers + body excerpt (300 chars)"
        except Exception as e:
            content, coverage = f"[EML error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext == ".msg":
        return "[MSG binary — Outlook message, filename signals only]", None, "binary, no extraction"

    elif ext == ".zip":
        try:
            import zipfile, tempfile
            with zipfile.ZipFile(file_path) as z:
                names = z.namelist()
            exts        = Counter(Path(n).suffix.lower() for n in names if Path(n).suffix)
            ext_summary = ", ".join(f"{e}×{c}" for e, c in exts.most_common(5))
            top_names   = ", ".join(names[:8])
            header      = f"[ZIP: {len(names)} files | types: {ext_summary} | samples: {top_names}]"
            PRIORITY_EXTS = [".pdf", ".docx", ".txt", ".xlsx", ".pptx", ".csv",
                             ".json", ".ifc", ".dxf", ".eml"]
            candidates = [n for n in names if Path(n).suffix.lower() in PRIORITY_EXTS
                          and not n.endswith("/") and Path(n).suffix.lower() != ".zip"]
            candidates.sort(key=lambda n: PRIORITY_EXTS.index(Path(n).suffix.lower()))
            to_sample     = candidates[:2]
            sampled_parts = []
            if to_sample:
                with zipfile.ZipFile(file_path) as z:
                    with tempfile.TemporaryDirectory() as tmp:
                        for name in to_sample:
                            z.extract(name, tmp)
                            extracted   = Path(tmp) / name
                            sub_content, _, sub_cov = _extract_content(extracted)
                            sampled_parts.append(
                                f"[sampled: {Path(name).name} | {sub_cov}]\n{sub_content}")
            content  = header + ("\n\n" + "\n\n".join(sampled_parts) if sampled_parts else "")
            content  = content[:DEFAULT_MAX_CHARS * 3]
            coverage = (f"{len(names)} files in archive"
                        + (f" — {len(to_sample)} file(s) content-sampled" if to_sample
                           else " — no text files to sample"))
        except Exception as e:
            content, coverage = f"[ZIP error: {e}]", "extraction failed"
        return content, None, coverage

    elif ext in (".shp", ".geojson", ".kml", ".gpkg"):
        # NEW ── try geopandas for richer spatial metadata
        try:
            import geopandas as gpd
            gdf      = gpd.read_file(str(file_path), rows=5)
            geom_types = gdf.geometry.geom_type.value_counts().to_dict()
            crs_str    = str(gdf.crs) if gdf.crs else "unknown CRS"
            cols       = list(gdf.columns)
            content  = (f"[Spatial file — CRS: {crs_str}]\n"
                        f"Geometry types: {geom_types}\n"
                        f"Columns: {cols}\n"
                        f"Feature count (sample): {len(gdf)}")[:DEFAULT_MAX_CHARS]
            coverage = f"geopandas — {len(gdf)} features sampled"
        except ImportError:
            # geopandas not available — fall back to raw text read
            try:
                raw      = file_path.read_text(encoding="utf-8", errors="replace")
                content  = raw[:400]
                coverage = _cov_chars(raw, 400, ext.upper())
            except Exception:
                size_kb  = file_path.stat().st_size // 1024
                content  = f"[{ext.upper()} spatial file — {size_kb} KB — likely binary]"
                coverage = "binary, no extraction"
        except Exception as e:
            content, coverage = f"[Spatial error: {e}]", "extraction failed"
        return content, None, coverage

    return "", None, "unsupported format"


# ── Public API ────────────────────────────────────────────────────────────
def layer1_technical(file_path: Path) -> dict:
    ext                      = file_path.suffix.lower()
    content, pages, coverage = _extract_content(file_path)
    year                     = _extract_year(file_path.name, content, file_path)  # ← passes file_path now
    info_type                = INFO_TYPE_MAP.get(ext, "Unknown")
    size_kb                  = file_path.stat().st_size // 1024

    return {
        # ── Original fields (unchanged) ──────────────────────────────────
        "filename":            file_path.name,
        "format":              ext.lstrip(".").upper(),
        "information_type":    info_type,
        "year":                year,
        "page_count":          pages,
        "size_kb":             size_kb,
        "folder":              file_path.parent.name,        # kept for backwards compat
        "file_path":           str(file_path),
        "content_sample":      content,
        "extraction_coverage": coverage,

        # ── NEW fields ───────────────────────────────────────────────────
        "size_category":       _size_category(size_kb),
        "path_segments":       _extract_path_segments(file_path),  # full folder chain → huge signal
        "filename_signals":    _extract_filename_signals(file_path.name),  # discipline/status/version/drawing_no
    }
